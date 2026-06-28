"""Builds a branded board-package .pptx that mirrors the full board book AND
its on-screen look: navy cover, white editorial section pages with serif
titles, light bordered tables, white stat cards, navy dividers for break
sections. Pure python-pptx. The client sends the active section list plus
every section's data; this renders one or more slides per section, in order."""
import io
import base64
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

WW, HH, MX, CW = 13.33, 7.5, 0.62, 12.09
LEFT, RIGHT, CENTER = PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.CENTER
TOP, MID = MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE
# palette mirrors the app CSS tokens
WHITE = RGBColor(0xFF, 0xFF, 0xFF); SOFT = RGBColor(0x5C, 0x67, 0x73); FAINT = RGBColor(0x97, 0xA1, 0xAC)
PAPER = RGBColor(0xF4, 0xF6, 0xF8); LINE = RGBColor(0xE2, 0xE7, 0xEC)
REDBG = RGBColor(0xFB, 0xEC, 0xEA); RED = RGBColor(0xB2, 0x3A, 0x30); REDINK = RGBColor(0x7A, 0x2A, 0x22)
INK = RGBColor(0x13, 0x20, 0x2B); RING1 = RGBColor(0x35, 0x55, 0x70); RING2 = RGBColor(0x2C, 0x4A, 0x64)
LABEL = RGBColor(0x8A, 0x9A, 0xA8); SUBWHITE = RGBColor(0xC7, 0xD2, 0xDC); POS = RGBColor(0x1F, 0x7A, 0x4D)
NEG = RGBColor(0xB2, 0x3A, 0x30)  # mirrors --neg (variance red)
DRAFTBG = RGBColor(0xFC, 0xEF, 0xDD); DRAFTINK = RGBColor(0xA9, 0x74, 0x1A)  # mirrors .tagDraft


def _hex(h, default):
    try:
        h = (h or default).lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        d = default.lstrip("#")
        return RGBColor(int(d[0:2], 16), int(d[2:4], 16), int(d[4:6], 16))


def _usd(n):
    try: return "$" + format(int(round(float(n))), ",")
    except Exception: return str(n) if n is not None else ""


def _M(n):
    try: return "$" + ("%.1f" % (float(n) / 1e6)) + "M"
    except Exception: return str(n) if n is not None else ""


def _var(a, b):
    try:
        d = float(a) - float(b)
        return ("+" if d >= 0 else "\u2212") + _usd(abs(d))
    except Exception:
        return ""


def _commaint(n):
    try:
        if isinstance(n, (int, float)): return "{:,}".format(int(n))
        return str(n) if n is not None else ""
    except Exception:
        return str(n) if n is not None else ""


def _image_bytes(src):
    """Bytes for an image from a data: URI or http(s) URL, or None.
    SVG and any fetch/decoding failure return None so the caller can fall
    back to the generated mark — an image is never allowed to break the deck."""
    if not src or not isinstance(src, str):
        return None
    try:
        if src.startswith("data:"):
            header, _, b64 = src.partition(",")
            if "svg" in header.lower() or not b64:
                return None
            return base64.b64decode(b64)
        if src.startswith("http://") or src.startswith("https://"):
            if src.lower().split("?")[0].endswith(".svg"):
                return None
            req = urllib.request.Request(src, headers={"User-Agent": "cu-library-pptx"})
            with urllib.request.urlopen(req, timeout=6) as resp:
                if "svg" in (resp.headers.get("Content-Type", "") or "").lower():
                    return None
                return resp.read()
    except Exception:
        return None
    return None


def _strip_theme(sp):
    """Remove theme/scheme color refs and the style element from a shape so the
    explicit RGB set in spPr always wins (prevents stray template accents)."""
    el = sp._element
    spPr = el.find(qn("p:spPr"))
    if spPr is None: spPr = el
    for sc in spPr.findall(".//" + qn("a:schemeClr")):
        par = sc.getparent()
        if par is not None: par.remove(sc)
    style_el = el.find(qn("p:style"))
    if style_el is not None: el.remove(style_el)


def build(p):
    theme = p.get("theme") or {}
    NAVY = _hex(theme.get("primary"), "#16314E")
    ACCENT = _hex(theme.get("accent"), "#1F9ED4")
    ACCENT_SOFT = _hex(theme.get("accentSoft"), "#ECF2F7")
    POSc = _hex(theme.get("pos"), "#1F7A4D")
    NEGc = _hex(theme.get("neg"), "#B23A30")
    logo_bytes = _image_bytes(theme.get("logoUrl"))
    cover_bytes = _image_bytes(theme.get("coverImage"))
    name = (theme.get("name") or "Credit Union")
    period = p.get("period") or {}
    plabel = period.get("label") or ""
    # mirrors the preview .printFooter: "{name} \u00b7 {period} Board Package \u00b7 Confidential"
    foot = name + ((u"  \u00b7  " + plabel) if plabel else "") + u"  \u00b7  Board Package  \u00b7  Confidential"
    runmeta = ((plabel + u"  \u00b7  Board Report") if plabel else "Board Report").upper()

    prs = Presentation(); prs.slide_width = Inches(WW); prs.slide_height = Inches(HH)
    BLANK = prs.slide_layouts[6]
    PG = [1]
    BARID = [1000]  # unique shape ids for hand-built bar XML

    def slide(bg=WHITE):
        s = prs.slides.add_slide(BLANK)
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-0.06), Inches(-0.06), Inches(WW + 0.12), Inches(HH + 0.12))
        r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
        _strip_theme(r)
        return s

    def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, radius=None, lcolor=None, lwpt=1.0):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.shadow.inherit = False
        if lcolor is not None:
            sp.line.color.rgb = lcolor; sp.line.width = Pt(lwpt)
        else:
            sp.line.fill.background()
        _strip_theme(sp)  # explicit RGB must win over any template accent (e.g. red/orange)
        if radius is not None:
            try: sp.adjustments[0] = radius
            except Exception: pass
        return sp

    def hline(s, x, y, w, color, thick=0.012):
        rect(s, x, y, w, thick, color)

    def outline(s, x, y, w, h, color, wpt=1.25):
        sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.background(); sp.line.color.rgb = color; sp.line.width = Pt(wpt); sp.shadow.inherit = False
        _strip_theme(sp)
        return sp

    def txt(s, x, y, w, h, runs, align=LEFT, anchor=TOP, font="Calibri", leading=None, spacing=False):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        for i, (t, sz, c, b) in enumerate(runs):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            par.alignment = align
            if leading is not None: par.line_spacing = leading
            r = par.add_run(); r.text = str(t); r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = font
        return tb

    def footer(s):
        hline(s, MX, 6.96, CW, LINE)
        txt(s, MX, 7.04, 10.6, 0.3, [(foot, 9, FAINT, False)])
        txt(s, WW - MX - 1.8, 7.04, 1.8, 0.3, [("PAGE " + str(PG[0]), 9, FAINT, False)], align=RIGHT)

    def place_img(s, raw, x, y, w=None, h=None):
        try:
            return s.shapes.add_picture(io.BytesIO(raw), Inches(x), Inches(y),
                                        Inches(w) if w else None, Inches(h) if h else None)
        except Exception:
            return None

    def fill_alpha(shape, pct_opaque):
        # set fill opacity (pct_opaque 0..100) so a cover photo shows through a wash
        try:
            srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
            a = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct_opaque * 1000))})
            srgb.append(a)
        except Exception:
            pass

    initials = ("".join(w[0] for w in name.split()[:2]).upper() or "CU")

    def logo_mark(s, x, y, size):
        # real uploaded/extracted logo carries the brand itself — show it alone
        if logo_bytes and place_img(s, logo_bytes, x, y, h=size):
            return
        rect(s, x, y, size, size, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.22, lcolor=ACCENT, lwpt=1.5)
        txt(s, x, y, size, size, [(initials, size * 34, WHITE, True)], align=CENTER, anchor=MID, font="Georgia")

    def wordmark(s, x, y, color=INK, h=0.3):
        # name with the final word in accent, mirroring the preview logoName
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(7), Inches(h)); tf = tb.text_frame
        tf.word_wrap = False; tf.vertical_anchor = MID
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        par = tf.paragraphs[0]; par.alignment = LEFT
        words = name.split()
        for i, wd in enumerate(words):
            r = par.add_run(); r.text = (wd if i == 0 else " " + wd)
            r.font.size = Pt(13); r.font.bold = True; r.font.name = "Georgia"
            r.font.color.rgb = ACCENT if (len(words) > 1 and i == len(words) - 1) else color

    TAGS = {
        "sourced": (ACCENT_SOFT, NAVY, "From core export"),
        "draft": (DRAFTBG, DRAFTINK, "AI draft · review"),
        "secured": (LINE, FAINT, "• Member data · secured"),
    }

    def pill(s, x, y, kind):
        bg, fg, label = TAGS[kind]
        w = 0.05 * len(label) + 0.26
        rect(s, x, y, w, 0.26, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        txt(s, x, y, w, 0.26, [(label, 8, fg, True)], align=CENTER, anchor=MID)
        return w

    def runhead(s):
        if logo_bytes:
            logo_mark(s, MX, 0.13, 0.34)
        else:
            logo_mark(s, MX, 0.15, 0.3)
            wordmark(s, MX + 0.42, 0.15, color=INK, h=0.3)
        txt(s, WW - MX - 4.5, 0.17, 4.5, 0.24, [(runmeta, 8.5, FAINT, True)], align=RIGHT, anchor=MID)
        hline(s, MX, 0.5, CW, LINE, thick=0.01)

    # white editorial page with serif title (mirrors .vTitle / .vLead)
    def page(title, sub=None):
        PG[0] += 1
        s = slide(WHITE)
        runhead(s)
        txt(s, MX, 0.66, CW, 0.66, [(title, 29, INK, True)], font="Georgia")
        if sub: txt(s, MX, 1.30, CW, 0.34, [(sub, 12.5, SOFT, False)])
        footer(s)
        return s, (1.86 if sub else 1.60)

    # navy divider band page (mirrors .divider)
    def divider_page(title, sub=None):
        PG[0] += 1
        s = slide(WHITE)
        runhead(s)
        rect(s, MX, 0.66, CW, 1.18, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        txt(s, MX + 0.4, 0.66, CW - 0.8, 1.18, [(title, 28, WHITE, True)], font="Georgia", anchor=MID)
        footer(s)
        top = 2.08
        if sub:
            txt(s, MX, top, CW, 0.32, [(sub, 12.5, SOFT, False)]); top += 0.46
        return s, top

    # white stat card (mirrors hGrid Card / hLabel / hValue / hNote)
    def tile(s, x, y, w, h, label, value, note, vsize=23):
        rect(s, x, y, w, h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, lcolor=LINE, lwpt=1.0)
        txt(s, x + 0.26, y + 0.2, w - 0.5, 0.3, [(label, 11.5, SOFT, False)])
        txt(s, x + 0.24, y + 0.48, w - 0.46, 0.5, [(value, vsize, INK, True)], font="Calibri")
        if note: txt(s, x + 0.26, y + h - 0.36, w - 0.5, 0.32, [(note, 10.5, FAINT, False)])

    def tiles_grid(s, items, top, h=1.7, cols=2):
        gx = 0.26
        w = (CW - gx * (cols - 1)) / cols
        for i, it in enumerate(items):
            l, v, n = it[0], it[1], it[2]
            vs = it[3] if len(it) > 3 else 23
            tile(s, MX + (i % cols) * (w + gx), top + (i // cols) * (h + 0.26), w, h, l, v, n, vs)

    # book-style table: light uppercase header, hairline rows, total emphasis
    def btable(s, headers, rows, colw, top, rh=0.4):
        tot = sum(colw); widths = [c / tot * CW for c in colw]
        xs = []; x = MX
        for w in widths: xs.append(x); x += w
        for ci, htext in enumerate(headers):
            al = LEFT if ci == 0 else RIGHT
            txt(s, xs[ci] + (0.02 if ci == 0 else 0), top, widths[ci] - 0.04, rh - 0.08, [(str(htext).upper(), 9, FAINT, True)], align=al, anchor=MID)
        hline(s, MX, top + rh - 0.06, CW, LINE)
        y = top + rh
        # Pass 1: row fills + separators, so no text can ever sit beneath a fill (z-order safe)
        yy = y
        for row in rows:
            if row[1]:
                rect(s, MX, yy, CW, rh, PAPER)
                hline(s, MX, yy, CW, INK, thick=0.022)
            else:
                hline(s, MX, yy + rh - 0.012, CW, LINE, thick=0.01)
            yy += rh
        # Pass 2: cell text on top
        yy = y
        for row in rows:
            cells = row[0]; total = row[1]
            cellcolors = row[2] if len(row) > 2 else None
            for ci, val in enumerate(cells):
                al = LEFT if ci == 0 else RIGHT
                col = (cellcolors or {}).get(ci)
                if col is None: col = NAVY if total else INK
                txt(s, xs[ci] + (0.04 if ci == 0 else 0), yy, widths[ci] - 0.06, rh, [(val, 11.5, col, total)], align=al, anchor=MID)
            yy += rh
        return yy

    def listrow(s, x, y, w, left, right, lbold=False, lsize=11, rsize=9.5, lcolor=INK, rcolor=FAINT, metaw=1.9):
        txt(s, x, y, w - metaw, 0.3, [(left, lsize, lcolor, lbold)], anchor=MID)
        if right is not None:
            txt(s, x + w - metaw, y, metaw, 0.3, [(right, rsize, rcolor, False)], align=RIGHT, anchor=MID)
        hline(s, x, y + 0.34, w, LINE, thick=0.01)

    # mini card (mirrors .card + .miniTitle)
    def minicard(s, x, y, w, h, title, meta, tag=None, tagcolor=None, body=None, tagkind=None):
        rect(s, x, y, w, h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, lcolor=LINE, lwpt=1.0)
        txt(s, x + 0.22, y + 0.16, w - 0.4, 0.5, [(title, 12, INK, True)])
        if body:
            txt(s, x + 0.22, y + 0.5, w - 0.4, 0.45, [(body, 10, SOFT, False)], leading=1.15)
        if meta: txt(s, x + 0.22, y + h - 0.62, w - 0.4, 0.22, [(meta, 9.5, FAINT, False)])
        if tagkind: pill(s, x + 0.22, y + h - 0.36, tagkind)
        elif tag: txt(s, x + 0.22, y + h - 0.34, w - 0.4, 0.22, [(tag.upper(), 8, tagcolor or ACCENT, True)])

    def bars(s, x, y, w, h, title, data, color):
        txt(s, x, y, w, 0.3, [(title, 13, INK, True)], font="Georgia")
        top = y + 0.5; chart_h = h - 0.5 - 0.32; base = top + chart_h
        data = [(d[0], float(d[1])) for d in data if d and d[1] is not None]
        if not data: return
        mx = max(v for _, v in data) or 1; n = len(data); slot = w / n; bw = slot * 0.52
        hexc = str(color)  # 'RRGGBB' — written straight into srgbClr, no theme involved
        for i, (yr, v) in enumerate(data):
            bh = (v / mx) * chart_h * 0.9; bx = x + i * slot + (slot - bw) / 2; by = base - bh
            bhh = max(bh, 0.02)
            BARID[0] += 1; sid = BARID[0]
            sp_xml = (
                '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                '<p:nvSpPr><p:cNvPr id="%d" name="bar%d"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
                '<p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
                '<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 16000"/></a:avLst></a:prstGeom>'
                '<a:solidFill><a:srgbClr val="%s"/></a:solidFill><a:ln><a:noFill/></a:ln></p:spPr>'
                '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
            ) % (sid, sid, int(bx * 914400), int(by * 914400), int(bw * 914400), int(bhh * 914400), hexc)
            s.shapes._spTree.append(parse_xml(sp_xml))
            txt(s, x + i * slot, by - 0.25, slot, 0.22, [(_M(v), 8.5, SOFT, False)], align=CENTER)
            txt(s, x + i * slot, base + 0.05, slot, 0.22, [(str(yr), 9.5, FAINT, False)], align=CENTER)
        hline(s, x, base, w, LINE)

    def narrcard(s, x, y, w, h, tag, text, draft=False):
        rect(s, x, y, w, h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04, lcolor=LINE, lwpt=1.0)
        txt(s, x + 0.3, y + 0.22, w - 2.0, 0.26, [(tag.upper(), 9, ACCENT, True)])
        if draft: pill(s, x + w - 1.4, y + 0.18, "draft")
        body = txt(s, x + 0.3, y + 0.56, w - 0.6, min(h - 0.8, 4.0), [(text, 12.5, INK, False)], font="Georgia", leading=1.4)
        try: body.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink long narratives to stay in-card
        except Exception: pass

    # ============================ COVER ============================
    def cover():
        s = slide(NAVY)
        if cover_bytes:
            place_img(s, cover_bytes, -0.06, -0.06, w=WW + 0.12, h=HH + 0.12)
            wash = rect(s, -0.06, -0.06, WW + 0.12, HH + 0.12, NAVY)
            fill_alpha(wash, 62)  # navy wash so the title stays legible over the photo
        outline(s, WW - 3.2, -3.4, 7.4, 7.4, RING1)
        outline(s, WW - 1.7, -2.0, 5.0, 5.0, RING2)
        if logo_bytes:
            place_img(s, logo_bytes, MX + 0.3, 0.58, h=0.52)
        else:
            logo_mark(s, MX + 0.3, 0.6, 0.5)
            wordmark(s, MX + 0.92, 0.6, color=WHITE, h=0.5)
        txt(s, MX + 0.3, 1.5, 11, 0.4, [(name.upper(), 13, ACCENT, True)])
        txt(s, MX + 0.27, 2.7, 11.5, 1.4, [("Board Report", 58, WHITE, True)], font="Georgia")
        if plabel: txt(s, MX + 0.3, 4.12, 11, 0.4, [(plabel.upper(), 14, SUBWHITE, False)])
        hline(s, MX + 0.3, 5.15, 3.0, ACCENT, thick=0.03)
        txt(s, MX + 0.3, 5.5, 5, 0.3, [("MEETING", 10, LABEL, True)])
        txt(s, MX + 0.3, 5.83, 6, 0.4, [(period.get("meeting") or "", 14, WHITE, False)])
        txt(s, 7.0, 5.5, 5, 0.3, [("LOCATION", 10, LABEL, True)])
        txt(s, 7.0, 5.83, 6, 0.4, [(period.get("location") or "", 14, WHITE, False)])

    # ============================ AGENDA ============================
    def agenda():
        A = p.get("agenda") or {}
        s, top = page("Board Meeting Agenda", ((period.get("meeting") or "") + ("  \u00b7  " + period.get("location") if period.get("location") else "")) or None)

        def group(x, y, w, title, rows):
            txt(s, x, y, w, 0.3, [(title, 11, NAVY, True)]); y += 0.34
            for r in (rows or []):
                item = ("(V)  " if r.get("vote") else "") + str(r.get("item", ""))
                meta = (str(r.get("owner") or "")) + ((" \u00b7 p." + str(r.get("page"))) if r.get("page") else "")
                listrow(s, x, y, w, item, meta, metaw=1.85); y += 0.4
            return y + 0.16

        colw = (CW - 0.5) / 2
        group(MX, top, colw, "CONSENT AGENDA", A.get("consent"))
        y2 = group(MX + colw + 0.5, top, colw, "BOARD ACTION ITEMS", A.get("actions"))
        group(MX + colw + 0.5, y2, colw, "REPORTS & UPDATES", A.get("reports"))

    # ============================ CONSENT ============================
    def consent():
        s, top = divider_page("Consent Agenda", "Items recommended for approval as a block")
        items = (p.get("agenda") or {}).get("consent") or []
        secured = {"Loan Charge-offs & Adjustments", "SAR Report", "New Memberships Report"}
        cols = 3; gx = 0.26; gy = 0.24; w = (CW - gx * (cols - 1)) / cols; h = 1.15
        for i, r in enumerate(items[:9]):
            x = MX + (i % cols) * (w + gx); y = top + (i // cols) * (h + gy)
            sec = r.get("item") in secured
            meta = (str(r.get("owner") or "")) + ((" \u00b7 p." + str(r.get("page"))) if r.get("page") else "")
            minicard(s, x, y, w, h, str(r.get("item", "")), meta, tagkind=("secured" if sec else "sourced"))

    # ============================ ACTION ITEMS ============================
    def actions():
        s, top = divider_page("Board Action Items", "Policy and administrative recommendations requiring a vote")
        items = p.get("actionItems") or []
        y = top; h = 1.04; gap = 0.16
        for a in items[:4]:
            rect(s, MX, y, CW, h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, lcolor=LINE, lwpt=1.0)
            txt(s, MX + 0.3, y + 0.16, CW - 2.2, 0.4, [("(V)  " + str(a.get("policy", "")), 13.5, INK, True)], font="Georgia")
            meta = (str(a.get("owner") or "")) + ((" \u00b7 p." + str(a.get("page"))) if a.get("page") else "")
            txt(s, MX + CW - 2.0, y + 0.2, 1.8, 0.3, [(meta, 10, FAINT, False)], align=RIGHT)
            txt(s, MX + 0.3, y + 0.58, CW - 1.7, 0.4, [(str(a.get("summary", "")), 11.5, SOFT, False)], leading=1.2)
            pill(s, MX + CW - 1.4, y + h - 0.36, "sourced")
            y += h + gap

    # ============================ FINANCIAL (multi) ============================
    def financial():
        s, top = page("Financial Report", "Statement of financial condition " + (period.get("financialsAsOf") or ""))
        nar = p.get("narrative") or ""
        if nar:
            nh = min(4.8, max(1.8, 0.018 * len(nar)))  # size the card to the text, no dead space
            narrcard(s, MX, top, CW, nh, "Financial narrative \u00b7 AI draft", nar, draft=True)

        def varcolor(a, b):
            try: return POSc if float(a) >= float(b) else NEGc
            except Exception: return None

        # highlights tiles
        s, top = page("Financial Highlights")
        _raw_hl = p.get("headline") or []
        print("[highlights] incoming headline count=%d labels=%s" % (
            len(_raw_hl), [h.get("label") for h in _raw_hl]), flush=True)
        hl = _raw_hl[:4]
        tiles = [(h.get("label", ""), h.get("value", ""), h.get("note", "")) for h in hl
                 if (h.get("label") or h.get("value"))]
        if len(tiles) < 4:
            # supplement from the statements so we never fall back to stale/placeholder labels
            bs = [r for r in (p.get("balanceSheet") or [])]
            isr = [r for r in (p.get("incomeStatement") or [])]
            def find(rows, *keys):
                for k in keys:
                    for r in rows:
                        if k in str(r.get("label", "")).lower(): return r
                return None
            have = {t[0].strip().lower() for t in tiles}
            cand = []
            a = find(bs, "total asset", "asset")
            if a: cand.append(("Total assets", _usd(a.get("actual")), "from balance sheet"))
            ln = find(bs, "net loan", "loan")
            if ln: cand.append(("Net loans", _usd(ln.get("actual")), "from balance sheet"))
            sh = find(bs, "share", "deposit")
            if sh: cand.append(("Total shares", _usd(sh.get("actual")), "from balance sheet"))
            ni = find(isr, "net income", "net")
            if ni: cand.append(("Net income YTD", _usd(ni.get("actual")), "from income statement"))
            for c in cand:
                if len(tiles) >= 4: break
                if c[0].strip().lower() not in have:
                    tiles.append(c); have.add(c[0].strip().lower())
        if not tiles:
            tiles = [("No headline data", "\u2014", "headline / statements were empty in the payload")]
        print("[highlights] rendered tiles=%s" % [t[0] for t in tiles[:4]], flush=True)
        tiles_grid(s, tiles[:4], top, h=1.85)
        # balance sheet (Actual / Budget / Variance / Prior YE \u2014 mirrors the preview BalanceTable)
        s, top = page("Statement of Financial Condition")
        pill(s, MX, top - 0.02, "sourced"); top += 0.36
        bs_rows = []
        for r in (p.get("balanceSheet") or []):
            a, b = r.get("actual"), r.get("budget")
            colors = {}; vc = varcolor(a, b)
            if vc is not None: colors[3] = vc
            bs_rows.append(([r.get("label", ""), _usd(a), _usd(b), _var(a, b), _usd(r.get("prior"))], bool(r.get("total")), colors))
        btable(s, ["Line", "Actual", "Budget", "Variance", "Prior YE"], bs_rows, [3.7, 2.05, 2.05, 2.2, 2.0], top)
        # income
        s, top = page("Income Statement (YTD)")
        pill(s, MX, top - 0.02, "sourced"); top += 0.36
        is_rows = []
        for r in (p.get("incomeStatement") or []):
            a, b = r.get("actual"), r.get("budget")
            colors = {}; vc = varcolor(a, b)
            if vc is not None: colors[3] = vc
            if r.get("net"): colors[1] = POSc
            is_rows.append(([r.get("label", ""), _usd(a), _usd(b), _var(a, b)], bool(r.get("total")), colors))
        btable(s, ["Line (YTD)", "Actual", "Budget", "Variance"], is_rows, [4.4, 2.55, 2.55, 2.55], top)
        # spread
        sp = p.get("spread") or []
        if sp:
            s, top = page("Spread Analysis")
            rows_h = (len(sp) + 1) * 0.4
            rect(s, MX, top - 0.12, CW, rows_h + 0.5, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03, lcolor=LINE, lwpt=1.0)
            btable(s, ["Line", "YTD 2025", "2024", "2023"],
                   [([r.get("label", ""), ("%.2f%%" % r.get("y2025")) if r.get("y2025") is not None else "", ("%.2f%%" % r.get("y2024")) if r.get("y2024") is not None else "", ("%.2f%%" % r.get("y2023")) if r.get("y2023") is not None else ""], bool(r.get("total"))) for r in sp],
                   [5.2, 2.3, 2.3, 2.3], top + 0.12)
        # trends
        lt = p.get("loanTrend") or []; sht = p.get("shareTrend") or []
        if lt or sht:
            s, top = page("Five-Year Trends")
            half = (CW - 0.5) / 2
            if lt: bars(s, MX, top, half, 4.6, "Net loans (5-yr)", [(d.get("yr"), d.get("v")) for d in lt], ACCENT)
            if sht: bars(s, MX + half + 0.5, top, half, 4.6, "Shares (5-yr)", [(d.get("yr"), d.get("v")) for d in sht], NAVY)

    # ============================ LENDING ============================
    def lending():
        s, top = page("Lending & Collections")
        L = p.get("lending") or {}; act = L.get("activity") or []
        a0 = act[0] if len(act) > 0 else {}; a1 = act[1] if len(act) > 1 else {}
        nar = p.get("lendingNarrative") or ""
        if nar:
            adv = min(2.2, max(1.1, 0.022 * len(nar)))  # reserve room for the full paragraph
            txt(s, MX, top, CW, adv, [(nar, 12, SOFT, False)], font="Georgia", leading=1.4); top += adv
        dr = L.get("delinquencyRatio"); drt = L.get("delinquencyRatioTable")
        try: flagged = (dr is not None and drt is not None and abs(float(dr) - float(drt)) > 0.01)
        except Exception: flagged = False
        bottom = 6.8 - (0.9 if flagged else 0.0)  # keep tiles (+flag) above the footer
        th = max(1.2, min(1.35 if nar else 1.55, (bottom - top) / 2 - 0.26))  # tall enough for the 3-line note
        litems = [
            ("Consumer loans YTD", _M(a0.get("ytd")), (str(a0.get("goal")) + " to goal") if a0.get("goal") else ""),
            ("Mortgage loans YTD", _M(a1.get("ytd")), (str(a1.get("goal")) + " to goal") if a1.get("goal") else ""),
            ("Delinquency ratio", str(L.get("delinquencyRatio", "")) + "%", (str(L.get("over60Count")) + " loans over 60 days") if L.get("over60Count") is not None else ""),
            ("Recoveries YTD", _usd(L.get("recoveriesYTD")), ("charge-offs " + _usd(L.get("chargeOffTotal"))) if L.get("chargeOffTotal") is not None else ""),
        ]
        tiles_grid(s, litems, top, h=th)
        fy = top + 2 * (th + 0.26)
        if flagged:
            rect(s, MX, fy, CW, 0.78, REDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1, lcolor=RED, lwpt=1.0)
            txt(s, MX + 0.3, fy + 0.12, 4, 0.3, [("CONSISTENCY FLAG", 10, RED, True)])
            txt(s, MX + 0.3, fy + 0.38, CW - 0.6, 0.36, [(("Delinquency reads %s%% in the narrative vs %s%% in the by-product table. Reconcile before distribution." % (dr, drt)), 11.5, REDINK, False)])

    # ============================ OPERATIONS ============================
    def ops():
        s, top = page("Operations & Risk Management")
        O = p.get("ops") or {}; sb = O.get("sb1075") or {}
        tm = _commaint(O.get("totalMembers"))
        oitems = [
            ("New members", O.get("newMembers", ""), ((tm + " total") + ((" \u00b7 " + str(O.get("closures")) + " closures") if O.get("closures") is not None else "")) if tm else ""),
            ("Member chatbot", "~" + str(O.get("sunnyConvosPerWeek", "")) + "/wk", ("live " + str(O.get("sunnyWeeks")) + " weeks") if O.get("sunnyWeeks") is not None else ""),
            ("SB 1075 fee cap", "$" + str(sb.get("planned", "")), ("down from $" + str(sb.get("current", "")) + ((" \u00b7 " + str(sb.get("effective"))) if sb.get("effective") else "")) if sb.get("current") is not None else ""),
            ("Escalated reviews", O.get("fraudReviews", ""), "member detail secured"),
        ]
        tiles_grid(s, oitems, top, h=1.45)
        sec = p.get("securedSections") or []
        if sec:
            ly = top + 2 * (1.45 + 0.26) + 0.05
            txt(s, MX, ly, 8, 0.3, [("Secured member-level sections", 11, NAVY, True)]); ly += 0.38
            for sx in sec[:4]:
                meta = (str(sx.get("owner") or "")) + ((" \u00b7 p." + str(sx.get("page"))) if sx.get("page") else "")
                listrow(s, MX, ly, CW, str(sx.get("name", "")) + "  \u2014  " + str(sx.get("count", "")), meta, metaw=2.0); ly += 0.4

    # ============================ SUBSIDIARY ============================
    def emerald():
        s, top = page("Subsidiary Report", "Wholly-owned CUSO \u00b7 sample subsidiary line of business")
        E = p.get("emerald") or {}
        items = [("Active accounts", str(E.get("accounts", "")), ""), ("Cumulative balances", str(E.get("balances", "")), ""), ("Concentration limit", str(E.get("concentrationLimit", "")), "")]
        tiles_grid(s, [(l, v, n, 22) for (l, v, n) in items], top, h=1.6, cols=3)
        mix = E.get("mix") or []
        if mix:
            my = top + 1.6 + 0.34
            txt(s, MX, my, 8, 0.3, [("Accounts by type", 11, NAVY, True)]); my += 0.4
            colw = (CW - 0.5) / 2
            for i, row in enumerate(mix):
                k = row[0] if len(row) > 0 else ""; v = row[1] if len(row) > 1 else ""
                cx = MX + (i % 2) * (colw + 0.5); cy = my + (i // 2) * 0.4
                listrow(s, cx, cy, colw, str(k), str(v), metaw=1.0, lsize=11.5, rcolor=NAVY, rsize=11.5)

    # ============================ CEO ============================
    def ceo():
        s, top = page("CEO / President's Report")
        nar = p.get("ceoNarrative") or ""
        bl = p.get("ceoHighlights") or ["No CEO report highlights were provided."]
        if nar:
            narrcard(s, MX, top, CW, 1.5, "Drafted from submitted highlights", nar, draft=True); top += 2.1
        txt(s, MX, top, 8, 0.3, [("Submitted highlights", 11, NAVY, True)], font="Calibri")
        tb = s.shapes.add_textbox(Inches(MX), Inches(top + 0.4), Inches(CW), Inches(6.4 - top)); tf = tb.text_frame; tf.word_wrap = True
        for i, b in enumerate(bl[:9]):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            par.space_after = Pt(7); par.line_spacing = 1.05
            r = par.add_run(); r.text = "\u2022   " + str(b); r.font.size = Pt(12.5); r.font.color.rgb = INK; r.font.name = "Calibri"

    # ============================ CALENDAR ============================
    def calendar():
        s, top = page("Calendar of Events")
        cal = p.get("calendar") or []
        y = top
        for row in cal[:13]:
            d = row[0] if len(row) > 0 else ""; n = row[1] if len(row) > 1 else ""
            txt(s, MX, y, 1.5, 0.3, [(str(d), 11.5, NAVY, True)], anchor=MID)
            txt(s, MX + 1.6, y, CW - 1.6, 0.3, [(str(n), 11, INK, False)], anchor=MID)
            hline(s, MX, y + 0.34, CW, LINE, thick=0.01)
            y += 0.4

    # ============================ APPENDIX ============================
    def appendix():
        s, top = divider_page("Appendix")
        ex = p.get("execs") or []; ac = p.get("acronyms") or []
        colw = (CW - 0.5) / 2
        txt(s, MX, top, colw, 0.3, [("LEADERSHIP", 11, NAVY, True)])
        y = top + 0.4
        for row in ex[:8]:
            nm = row[0] if len(row) > 0 else ""; rl = row[1] if len(row) > 1 else ""
            txt(s, MX, y, colw, 0.28, [(str(nm), 11.5, INK, True)], anchor=MID)
            txt(s, MX, y + 0.26, colw, 0.24, [(str(rl), 9.5, SOFT, False)])
            y += 0.6
        x2 = MX + colw + 0.5
        txt(s, x2, top, colw, 0.3, [("ACRONYMS", 11, NAVY, True)])
        y = top + 0.4
        for row in ac[:12]:
            a = row[0] if len(row) > 0 else ""; d = row[1] if len(row) > 1 else ""
            txt(s, x2, y, 0.95, 0.26, [(str(a), 10.5, NAVY, True)], anchor=MID)
            txt(s, x2 + 1.0, y, colw - 1.0, 0.26, [(str(d), 10, INK, False)], anchor=MID)
            y += 0.36

    DISPATCH = {
        "agenda": agenda, "consent": consent, "actions": actions, "financial": financial,
        "lending": lending, "ops": ops, "emerald": emerald, "ceo": ceo,
        "calendar": calendar, "appendix": appendix,
    }
    DEFAULT = ["agenda", "consent", "actions", "financial", "lending", "ops", "emerald", "ceo", "calendar", "appendix"]

    cover()
    sections = p.get("sections") or DEFAULT
    seen = set()
    for sid in sections:
        if sid == "cover" or sid in seen:
            continue
        seen.add(sid)
        fn = DISPATCH.get(sid)
        if fn:
            try: fn()
            except Exception:
                pass

    out = io.BytesIO(); prs.save(out); return out.getvalue()
